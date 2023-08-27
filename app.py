from pptx import Presentation
from pptx.util import Inches

# テキストファイルからスライドデータを読み込む
with open("data/slide-data.txt", "r", encoding="utf-8") as f:
    slides_data = f.read().splitlines()

# プレゼンテーションの作成とスライドの生成
prs = Presentation()

# リストの長さを2で割った回数までループ
for i in range(0, len(slides_data)//2):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # リストからタイトルとコンテンツを取得
    title = slides_data[i * 2]
    content = slides_data[i * 2 + 1]

    title_box = slide.shapes.title
    title_box.text = title

    content_box = slide.placeholders[1]
    content_box.text = content

# ファイル保存
prs.save("高校野球と熱中症対策プレゼンテーション.pptx")
